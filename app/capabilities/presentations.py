from __future__ import annotations
class PresentationCapability:
    def inspect(self,path:str)->dict:
        from pptx import Presentation
        p=Presentation(path); return {"slides":len(p.slides),"titles":[s.shapes.title.text if s.shapes.title else None for s in p.slides]}
